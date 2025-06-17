from pathlib import Path
from typing import List
import dtlpy as dl
import tempfile
import logging
import os
from pptx import Presentation

logger = logging.getLogger('ppt-to-text-logger')


class PptExtractor(dl.BaseServiceRunner):

    def ppt_extraction(self, item: dl.Item, context: dl.Context) -> List[dl.Item]:
        """
        Extract text from PPT/PPTX item and upload it as text files.

        :param item: Dataloop item, PPT/PPTX file
        :param context: Dataloop context containing node configuration
        :return: List of Text Dataloop items
        """
        node = context.node
        extract_notes = node.metadata['customNodeConfig']['extract_notes']
        extract_tables = node.metadata['customNodeConfig']['extract_tables']
        remote_path_for_extractions = node.metadata['customNodeConfig']['remote_path_for_extractions']

        # Check if the item is a PowerPoint file
        suffix = Path(item.name).suffix.lower()
        if suffix not in {'.ppt', '.pptx'}:
            raise ValueError(f"Item id: {item.id} is not a PowerPoint file! This function expects .ppt or .pptx only")

        # Download item
        with tempfile.TemporaryDirectory() as temp_dir:
            item_local_path = item.download(local_path=temp_dir)

            new_items_path = self.extract_text_from_ppt(
                ppt_path=item_local_path,
                extract_notes=extract_notes,
                extract_tables=extract_tables
            )

            new_items = item.dataset.items.upload(
                local_path=new_items_path,
                remote_path=remote_path_for_extractions,
                item_metadata={
                    'user': {
                        'extracted_from_ppt': True,
                        'original_item_id': item.id
                    }
                },
                overwrite=True
            )

            if new_items is None:
                raise dl.PlatformException(f"No items were uploaded! local paths: {new_items_path}")
            elif isinstance(new_items, dl.Item):
                all_items = [new_items]
            else:
                all_items = [item for item in new_items]

        return all_items

    @staticmethod
    def extract_text_from_ppt(ppt_path: str, extract_notes: bool = False, extract_tables: bool = False) -> List[str]:
        """
        Extracts text from a PowerPoint file and saves each slide as a separate .txt file.

        Args:
            ppt_path (str): The path to the PowerPoint file to be processed.
            extract_notes (bool): Whether to extract speaker notes from slides.
            extract_tables (bool): Whether to extract text from tables in slides.

        Returns:
            list: A list of paths to the generated .txt files, each corresponding to a slide in the presentation.
        """
        presentation = Presentation(ppt_path)
        text_files = []

        # Loop through each slide and save text to individual .txt files
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from tables if requested
                if extract_tables and shape.shape_type == 19:  # Table shape type
                    table_text = PptExtractor._extract_table_text(shape.table)
                    if table_text:
                        slide_text.append(table_text)
            
            # Extract speaker notes if requested
            if extract_notes and slide.notes_slide.notes_text_frame.text.strip():
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_text.append(f"\n--- Speaker Notes ---\n{notes_text}")

            # Combine all text from the slide
            combined_text = '\n'.join(slide_text)
            
            # Define the output path for each slide's text
            new_item_path = f'{os.path.splitext(ppt_path)[0]}_slide_{slide_num + 1}.txt'
            
            with open(new_item_path, 'w', encoding='utf-8') as f:
                f.write(combined_text)

            # Add the file path to the list of text files
            text_files.append(new_item_path)
            logger.info(f"Extracted text from slide {slide_num + 1}, saved to {new_item_path}")

        return text_files

    @staticmethod
    def _extract_table_text(table) -> str:
        """
        Extract text from a PowerPoint table.

        Args:
            table: PowerPoint table object

        Returns:
            str: Formatted table text with rows separated by newlines and cells by tabs
        """
        table_text = []
        
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_text.append(cell_text)
            table_text.append('\t'.join(row_text))
        
        return '\n'.join(table_text) 