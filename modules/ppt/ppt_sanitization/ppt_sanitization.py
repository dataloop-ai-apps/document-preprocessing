import openai
import os
import cv2
import base64
import shutil
import pptx
import logging
import dtlpy as dl
import numpy as np
from pptx import Presentation
from typing import Union

logger = logging.getLogger("[PPT-Sanitization]")

NER_PROMPT_MESSAGE = """
In the following text, follow these guidelines to respond with the corrected text. Only send back the corrected text, do not decorate with any other text:
1. If the text is a name of a company, organization, non-profit institute of any type, replace the name with [Org]
2. If the text is a location such as name of a country, city, village, street, state, site or region, replace it with [Location]
3. If the text is a person name, replace with [Person]
4. If the text is a project name in the structure of 'Project Name', replace this text with [Project]
5. Replace all currency symbols and currency codes (like USD, SAR, EURO, etc.) with the fixed text '[Currency]'.
   Ensure that other aspects of the text remain unchanged. Here are some examples for your reference:
   * Original: 'I spent $100 on groceries.'/ Rewritten: 'I spent [Currency]100 on groceries.
   * Original: 'The price was 150 EUR for the tickets.'/ Rewritten: 'The price was 150 [Currency] for the tickets.'
6. If the text is numerical financial data in any currency, keep the currency symbol and replace the numerical data with [xx]
7. If the text is a stock ticker, replace it with [Stock] We have a pragmatic and proven approach to drive efficiency (IT cost re-balancing) while managing effectiveness (productivity, quality, speed) to unlock business value. We call it project Blackwell. Robert Lewangoalski from Croatia is leading our finance division, responsible for SAR   €140M annual sales . We partner with mount sinai hospital from Norway (NASDAQ: AMZN) for student exchange and medical collaboration     
"""

VISUAL_IDENTITY_PROMPT_MESSAGE = """
Consider the following definition of what a logo is: 
'A logo (abbreviation of logotype) is a graphic mark, emblem, or symbol used to aid and promote public identification and recognition. It may be of an abstract or figurative design or to include the text of the name that it represents as in a wordmark.'
Also, consider this distinction: an informational image contains graphs, plots, or other sorts of information.
a decorative image are shapes, colors, characters, or other images that a company would use to convey their identity in a presentation, but not business information.
Now, answer these two questions: does the image contain a logo? is the image informational or decorative?
The first question should be answered solely with yes or no. The second question should be answered solely with either informational or decorative. 
That is, answer two words, separated by a comma, each one answering one of the questions.
"""


class RemoveSensitiveText(dl.BaseServiceRunner):
    def __init__(self, openai_key):
        self.client = openai.OpenAI(api_key=os.environ.get(openai_key))
        self.ner_prompt_message = NER_PROMPT_MESSAGE
        self.visual_identity_prompt_message = VISUAL_IDENTITY_PROMPT_MESSAGE

    @staticmethod
    def copy_text_attributes(source_run, target_run):
        # Copy font type, size, bold, italic, underline, color
        target_run.font.bold = source_run.font.bold
        target_run.font.italic = source_run.font.italic
        target_run.font.underline = source_run.font.underline
        target_run.font.size = source_run.font.size

    @staticmethod
    def chatgpt_request(content: Union[str, list],
                        system_prompt: str,
                        max_tokens: int = 10,
                        gpt_model: str = "gpt-3.5-turbo"):
        response = openai.chat.completions.create(
            model=gpt_model,
            messages=[
                {"role": "system",
                 "content": system_prompt},
                {"role": "user",
                 "content": content}
                ],
            max_tokens=max_tokens
            )
        return response.choices[0].message.content

    @staticmethod
    def create_image_content_gpt(image_buffer):
        encoded_image = base64.b64encode(image_buffer).decode('utf-8')
        content = [
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{encoded_image}"
                    }
                }
            ]
        return content

    @staticmethod
    def get_shape_coords(shape):
        return [shape.left, shape.top, shape.width, shape.height]

    def clean_images(self, shape, new_slide):
        if shape.shape_type == 13:  # image
            image = shape.image
            identification_answer = self.chatgpt_request(
                content=self.create_image_content_gpt(image.blob),
                system_prompt=self.visual_identity_prompt_message,
                gpt_model="gpt-4o",
                max_tokens=10
                ).lower()
            is_visual_identity_element = 'yes' in identification_answer or 'decorative' in identification_answer
            if is_visual_identity_element:
                with open(image.filename, 'bw') as f:
                    print(f"Logo was present, generating black image of size {image.width}x{image.height}")
                    rectangle_image = np.zeros((image.height, image.width, 3), dtype=np.uint8)
                    color = (0, 0, 0)  # Black color
                    thickness = -1  # Fill the rectangle
                    cv2.rectangle(rectangle_image, (0, 0), (image.width, image.height), color, thickness)
                    cv2.imwrite(image.filename, rectangle_image)
            else:
                print("Logo was not present. Image will be kept")
                with open(image.filename, 'bw') as f:
                    f.write(image.blob)
            new_slide.shapes.add_picture(image.filename, *self.get_shape_coords(shape))
        elif hasattr(shape, "text"):
            text_box = new_slide.shapes.add_textbox(*self.get_shape_coords(shape))
            text_frame = text_box.text_frame
            text_frame.text = shape.text
        else:
            logger.info("Shape is neither image nor next. Ignoring.")

    def clean_text(self, shape, new_slide):
        if hasattr(shape, "text"):
            text_frame = shape.text_frame
            tx_box = new_slide.shapes.add_textbox(*self.get_shape_coords(shape))
            tf = tx_box.text_frame
            tf.text = ""  # Clear the default text
            # Copy text and formatting
            for paragraph in text_frame.paragraphs:
                p = tf.add_paragraph()
                p.alignment = paragraph.alignment
                for run in paragraph.runs:
                    new_run = p.add_run()
                    response = self.chatgpt_request(run.text,
                                                    self.ner_prompt_message,
                                                    max_tokens=int(1.5 * len(run.text)),
                                                    gpt_model="chatgpt-3.5-turbo")
                    print(
                        f"old text: {run.text}, @@@@ new text: {response}")
                    new_run.text = response
                    self.copy_text_attributes(run, new_run)
        elif shape.shape_type == 13:
            new_slide.shapes.add_picture(shape.image.blob, *self.get_shape_coords(shape))
        else:
            logger.info("Shape is neither image nor next. Ignoring.")

    def sanitize(self, item: dl.Item, element: str):
        # Load the existing presentation
        data_dir = os.path.join(os.getcwd(), 'data')
        os.makedirs(data_dir, exist_ok=True)
        path = item.download(local_path=data_dir)
        prs = pptx.Presentation(path)
        # Create a new presentation object to hold the cleaned data
        new_prs = Presentation()
        # Iterate through all the slides in the original presentation
        for slide in prs.slides:
            # Create a blank slide in the new presentation
            blank_slide_layout = new_prs.slide_layouts[6]  # Using title only layout for more space
            new_slide = new_prs.slides.add_slide(blank_slide_layout)

            # Iterate through all shapes in the slide
            for shape in slide.shapes:
                if element == "image":
                    self.clean_images(shape, new_slide)
                elif element == "text":
                    self.clean_text(shape, new_slide)

        # Save the new presentation
        new_item_path = os.path.join(data_dir, f'{element}_sanitized_{item.name}')
        new_prs.save(new_item_path)
        new_item = item.dataset.items.upload(
            local_path=new_item_path,
            remote_path="/no_theme",
            overwrite=True,
            raise_on_error=True,
        )
        shutil.rmtree(data_dir)
        return new_item

    def sanitize_text(self, item: dl.Item):
        return self.sanitize(item, "text")

    def sanitize_visual_identity(self, item: dl.Item):
        return self.sanitize(item, "image")
